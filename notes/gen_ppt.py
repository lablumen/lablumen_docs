#!/usr/bin/env python3
"""Generate lablumen-presentation.pptx — LabLumen DevOps Capstone Deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
import os

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(10, 18, 52)
TEAL   = RGBColor(0, 172, 162)
TEAL_D = RGBColor(0, 110, 107)
LGREY  = RGBColor(205, 215, 232)
ORANGE = RGBColor(255, 145, 20)
GREEN  = RGBColor(50, 200, 90)
WHITE  = RGBColor(255, 255, 255)
CARD   = RGBColor(16, 28, 68)
CODEBG = RGBColor(12, 22, 48)
LTEAL  = RGBColor(120, 225, 215)
YELLOW = RGBColor(255, 215, 0)

def H(r, g, b): return f"{r:02X}{g:02X}{b:02X}"

NH = H(10,18,52); TH = H(0,172,162); TDH = H(0,110,107)
CH = H(16,28,68); CBH = H(12,22,48); BH = H(0,0,0)
OH = H(255,145,20); GH = H(50,200,90)

SW = Inches(13.333); SH = Inches(7.5)
CY = Inches(1.18);   CEH = Inches(6.0)
CL = Inches(0.35);   CW = Inches(12.6)
C1W = Inches(5.95);  C2L = Inches(6.65); C2W = Inches(6.35)

_sid = [500]

# ── Core helpers ─────────────────────────────────────────────────────────────
def rect(sl, l, t, w, hh, fill, bdr=None, bw=1.5):
    _sid[0] += 1; sid = _sid[0]
    ln = (f'<a:ln w="{int(Pt(bw))}"><a:solidFill><a:srgbClr val="{bdr}"/>'
          f'</a:solidFill></a:ln>') if bdr else '<a:ln><a:noFill/></a:ln>'
    xml = (f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
           f'<p:nvSpPr><p:cNvPr id="{sid}" name="s{sid}"/>'
           f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
           f'<p:spPr><a:xfrm><a:off x="{int(l)}" y="{int(t)}"/>'
           f'<a:ext cx="{int(w)}" cy="{int(hh)}"/></a:xfrm>'
           f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
           f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>{ln}'
           f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>')
    sl.shapes._spTree.append(parse_xml(xml.encode()))

def set_bg(sl, r, g, b):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = RGBColor(r, g, b)

def ns(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def tx(sl, text, l, t, w, hh, sz=12, bold=False, col=None,
       align=PP_ALIGN.LEFT, italic=False, fn="Calibri"):
    box = sl.shapes.add_textbox(l, t, w, hh)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.name = fn
    if col: run.font.color.rgb = col
    return box

def ml(sl, lines, l, t, w, hh, sz=11, col=None, spc=3,
       align=PP_ALIGN.LEFT, fn="Calibri", bold_set=None):
    """Multi-line textbox. lines: list of str|(str,color) tuples."""
    col = col or LGREY; bold_set = bold_set or set()
    box = sl.shapes.add_textbox(l, t, w, hh)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        text, c = (item if isinstance(item, tuple) else (item, col))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(spc)
        run = p.add_run(); run.text = text
        run.font.size = Pt(sz); run.font.name = fn
        run.font.bold = (i in bold_set); run.font.color.rgb = c
    return box

def hbar(sl, title, sub=None):
    rect(sl, 0, 0, SW, Inches(1.12), TH)
    rect(sl, 0, Inches(1.12), SW, Inches(0.04), TDH)
    tx(sl, title, Inches(0.35), Inches(0.1), Inches(12.5), Inches(0.8),
       sz=27, bold=True, col=WHITE)
    if sub:
        tx(sl, sub, Inches(0.35), Inches(0.83), Inches(12.5), Inches(0.32),
           sz=10, col=RGBColor(195, 238, 234))
    rect(sl, 0, SH-Inches(0.27), SW, Inches(0.27), H(8,14,42))
    tx(sl, "LabLumen  |  Cloud Native Microservices Platform  |  DevOps Capstone 2026",
       Inches(0.35), SH-Inches(0.25), Inches(11), Inches(0.25),
       sz=7.5, col=RGBColor(70,95,145))

def sec_slide(prs, num, title, sub=""):
    sl = ns(prs); set_bg(sl, 0, 110, 107)
    rect(sl, 0, SH-Inches(0.07), SW, Inches(0.07), OH)
    tx(sl, num, Inches(0.5), Inches(0.9), Inches(3.5), Inches(3.8),
       sz=110, bold=True, col=RGBColor(0,90,87))
    tx(sl, title, Inches(0.55), Inches(2.7), Inches(11.5), Inches(2.2),
       sz=42, bold=True, col=WHITE)
    if sub:
        tx(sl, sub, Inches(0.55), Inches(4.65), Inches(11), Inches(0.6),
           sz=15, col=RGBColor(195, 238, 234))
    return sl

def ph_slide(prs, label, sub="Replace with actual screenshot"):
    sl = ns(prs); set_bg(sl, 0, 0, 0)
    rect(sl, Inches(0.28), Inches(0.28), SW-Inches(0.56), SH-Inches(0.56),
         BH, TH, bw=2.5)
    tx(sl, "[ Screenshot / Diagram Placeholder ]",
       Inches(1), Inches(1.4), SW-Inches(2), Inches(0.7),
       sz=13, col=RGBColor(80,80,80), align=PP_ALIGN.CENTER, bold=True)
    tx(sl, label, Inches(0.6), Inches(2.5), SW-Inches(1.2), Inches(2),
       sz=32, bold=True, col=TEAL, align=PP_ALIGN.CENTER)
    tx(sl, sub, Inches(1.5), Inches(5.0), SW-Inches(3), Inches(0.7),
       sz=12, col=RGBColor(100,100,100), align=PP_ALIGN.CENTER, italic=True)
    return sl

def card(sl, l, t, w, hh, title=None):
    rect(sl, l, t, w, hh, CH, TH, bw=0.75)
    if title:
        rect(sl, l, t, w, Inches(0.035), TH)
        tx(sl, title, l+Inches(0.12), t+Inches(0.08), w-Inches(0.24), Inches(0.26),
           sz=10, bold=True, col=TEAL)

def codebox(sl, lines, l, t, w, hh, title=None):
    rect(sl, l, t, w, hh, CBH, TH, bw=0.5)
    yo = Inches(0.06)
    if title:
        rect(sl, l, t, w, Inches(0.26), TH)
        tx(sl, title, l+Inches(0.1), t+Inches(0.04), w-Inches(0.2), Inches(0.22),
           sz=8, bold=True, col=WHITE)
        yo = Inches(0.3)
    box = sl.shapes.add_textbox(l+Inches(0.12), t+yo, w-Inches(0.24), hh-yo-Inches(0.06))
    tf = box.text_frame; tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        run = p.add_run(); run.text = line
        run.font.size = Pt(8.5); run.font.name = "Courier New"
        run.font.color.rgb = LTEAL

def tbl(sl, data, l, t, w, hh, col_ws):
    """Dark-themed table. data[0] = header row."""
    rows = len(data); cols = len(data[0])
    table = sl.shapes.add_table(rows, cols, l, t, w, hh).table
    for ci, cw in enumerate(col_ws):
        table.columns[ci].width = cw
    for ri, row in enumerate(data):
        is_hdr = ri == 0
        for ci, cell_txt in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            if is_hdr:
                cell.fill.fore_color.rgb = TEAL
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = CARD
            else:
                cell.fill.fore_color.rgb = NAVY
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs: run.text = ""
            run = p.add_run(); run.text = cell_txt
            run.font.size = Pt(10 if is_hdr else 9.5)
            run.font.bold = is_hdr; run.font.name = "Calibri"
            run.font.color.rgb = WHITE if is_hdr else LGREY

def chk(text, color=None):
    return (f"✔  {text}", color or GREEN)

def bull(text, color=None):
    return (f"▶  {text}", color or LGREY)

def sub(text, color=None):
    return (f"    •  {text}", color or RGBColor(170,185,210))

# ═══════════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    # ── 1. TITLE ──────────────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    rect(sl, 0, Inches(2.3), SW, Inches(3.0), H(14,26,70))
    rect(sl, 0, Inches(2.3), Inches(0.08), Inches(3.0), TH)
    rect(sl, 0, SH-Inches(0.07), SW, Inches(0.07), OH)
    tx(sl, "LabLumen", Inches(1), Inches(2.45), Inches(11), Inches(1.5),
       sz=72, bold=True, col=WHITE)
    tx(sl, "Cloud Native Microservices Platform",
       Inches(1), Inches(3.9), Inches(11), Inches(0.65),
       sz=24, col=TEAL)
    tx(sl, "AWS EKS   |   Terraform   |   ArgoCD   |   GitHub Actions",
       Inches(1), Inches(4.55), Inches(11), Inches(0.5),
       sz=14, col=LGREY)
    rect(sl, Inches(1), Inches(5.25), Inches(4.5), Inches(0.03), TH)
    tx(sl, "Raphel M L   |   DevOps Capstone Review   |   June 2026",
       Inches(1), Inches(5.35), Inches(11), Inches(0.45),
       sz=12, col=RGBColor(130,150,190), italic=True)

    # ── 2. AGENDA ─────────────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Agenda")
    sections = [
        ("01", "The Application", "Services, Tech Stack, Event-Driven Architecture"),
        ("02", "Cloud Infrastructure", "AWS Services, Architecture Overview"),
        ("03", "Terraform — IaC", "Modules, State Management, Pipelines"),
        ("04", "Kubernetes on EKS", "Namespaces, ArgoCD, Karpenter, Grafana"),
        ("05", "GitHub Actions — CI/CD", "Branching Strategy, Security Gates, Pipelines"),
        ("06", "Security & Results", "Zero static credentials, what was achieved"),
    ]
    for i, (num, title, desc) in enumerate(sections):
        col = i % 2
        row = i // 2
        l = CL + col * Inches(6.45)
        t = CY + Inches(0.1) + row * Inches(1.85)
        card(sl, l, t, Inches(6.1), Inches(1.65))
        tx(sl, num, l+Inches(0.15), t+Inches(0.12), Inches(0.9), Inches(0.8),
           sz=30, bold=True, col=TEAL)
        tx(sl, title, l+Inches(1.05), t+Inches(0.15), Inches(4.8), Inches(0.5),
           sz=14, bold=True, col=WHITE)
        tx(sl, desc, l+Inches(1.05), t+Inches(0.65), Inches(4.8), Inches(0.7),
           sz=9.5, col=LGREY)

    # ── 3. SECTION: THE APPLICATION ───────────────────────────────────────────
    sec_slide(prs, "01", "The Application",
              "Microservices  |  Tech Stack  |  Event-Driven Architecture")

    # ── 4. PROJECT OVERVIEW ───────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Project Overview — What is LabLumen?")
    ml(sl, [
        bull("Cloud-native laboratory platform — patients book lab tests, staff manage workflows, AI generates report summaries"),
        bull("Built as independent microservices following 12-factor app principles"),
        bull("Full DevOps lifecycle:  IaC provisioning  →  containerised deploy  →  GitOps CD  →  observability"),
        bull("Two environments on a single EKS cluster:"),
        sub("api-dev.rnld101.xyz   →   lablumen-dev namespace  (dev)"),
        sub("api.rnld101.xyz          →   lablumen namespace      (prod)"),
        bull("Auth: AWS Cognito (SRP flow) — separate patient & staff roles"),
        bull("Frontend SPA hosted on S3 + CloudFront at app.rnld101.xyz"),
    ], CL, CY+Inches(0.1), C1W+Inches(0.2), Inches(5.5), sz=11.5, spc=6)

    card(sl, C2L, CY+Inches(0.1), C2W, Inches(2.8), "Capstone Evaluation Pillars")
    ml(sl, [
        chk("Infrastructure as Code (Terraform)"),
        chk("Kubernetes Deployment (EKS)"),
        chk("CI/CD Pipelines (GitHub Actions)"),
        chk("Cloud Integration (IRSA, no static keys)"),
    ], C2L+Inches(0.15), CY+Inches(0.55), C2W-Inches(0.3), Inches(2.1), sz=11, spc=7)

    card(sl, C2L, CY+Inches(3.05), C2W, Inches(2.85), "Bonus Pillars Achieved")
    ml(sl, [
        (f"★  ArgoCD GitOps Deployment", ORANGE),
        (f"★  AI/ML Integration (Bedrock + Textract)", ORANGE),
        (f"★  Prometheus + Grafana Monitoring", ORANGE),
        (f"★  Event-Driven Architecture (SQS)", ORANGE),
        (f"★  Infracost Cost Estimation", ORANGE),
    ], C2L+Inches(0.15), CY+Inches(3.5), C2W-Inches(0.3), Inches(2.2), sz=11, spc=7)

    # ── 5. TECH STACK ─────────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Technology Stack")

    card(sl, CL, CY+Inches(0.05), C1W, Inches(2.75), "Backend Services")
    ml(sl, [
        bull("Python 3.12 + FastAPI  (async, 3 backend services)"),
        bull("SQLAlchemy async + Alembic  (DB migrations)"),
        bull("PostgreSQL RDS + Redis in-cluster"),
        bull("AWS SQS  (event bus)  |  SES  (email)"),
        bull("AWS Bedrock Nova Lite + Textract  (AI pipeline)"),
    ], CL+Inches(0.15), CY+Inches(0.5), C1W-Inches(0.3), Inches(2.1), sz=11, spc=5)

    card(sl, CL, CY+Inches(3.0), C1W, Inches(2.85), "Frontend & Auth")
    ml(sl, [
        bull("React + TypeScript + Vite  SPA"),
        bull("AWS Cognito  —  patient & staff roles, SRP auth"),
        bull("S3 + CloudFront hosting  (HTTPS, OAC)"),
        bull("Vite build deployed via GitHub Actions OIDC"),
    ], CL+Inches(0.15), CY+Inches(3.5), C1W-Inches(0.3), Inches(2.1), sz=11, spc=5)

    card(sl, C2L, CY+Inches(0.05), C2W, Inches(2.75), "Platform & DevOps")
    ml(sl, [
        bull("Terraform  —  all AWS infra as code  (12 modules)"),
        bull("Docker  (multistage builds, non-root)  +  ECR registry"),
        bull("AWS EKS  (Kubernetes v1.31)  +  Karpenter autoscaler"),
        bull("ArgoCD  (GitOps, App-of-Apps, 16 apps)"),
        bull("GitHub Actions  —  reusable polyrepo CI/CD engine"),
    ], C2L+Inches(0.15), CY+Inches(0.5), C2W-Inches(0.3), Inches(2.1), sz=11, spc=5)

    card(sl, C2L, CY+Inches(3.0), C2W, Inches(2.85), "Security & Observability")
    ml(sl, [
        bull("GitHub OIDC  +  IRSA  —  zero static credentials"),
        bull("SonarCloud SAST  |  Snyk SCA  |  Trivy container scan"),
        bull("Checkov IaC scan  (SARIF  →  GitHub Security tab)"),
        bull("Prometheus + Grafana  (kube-prometheus-stack)"),
        bull("KMS CMK  encrypts ECR images + Secrets Manager"),
    ], C2L+Inches(0.15), CY+Inches(3.5), C2W-Inches(0.3), Inches(2.1), sz=11, spc=5)

    # ── 6. MICROSERVICES TABLE ────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Microservices Breakdown")
    rows = [
        ["Service", "Runtime", "Role & Responsibilities"],
        ["appointment-service", "FastAPI / Python 3.12",
         "Booking, patients, lab-test catalogue. Redis slot-locks prevent double-booking."
         " Publishes appointment.booked events to SQS. Alembic DB migrations on startup."],
        ["report-service", "FastAPI / Python 3.12",
         "Lab report upload (S3) and delivery (pre-signed URLs). AI-powered RAG chat"
         " via Bedrock Nova Lite + pgvector. IRSA-authenticated for S3 & Bedrock access."],
        ["notification-service", "FastAPI worker / Python 3.12",
         "Long-polls SQS for appointment events. Sends confirmation emails via SES"
         " (domain identity, DKIM). No public ingress. IRSA for SQS + SES access."],
        ["ai-service", "AWS Lambda / SAM",
         "S3-triggered pipeline: Textract OCR → Bedrock summarise → Titan embed → pgvector."
         " Cross-account Bedrock role assumed at cold start. NOT deployed in Kubernetes."],
        ["frontend", "React + TypeScript + Vite",
         "Patient portal (book tests, view/chat reports) and staff portal (manage queue,"
         " update status). AWS Cognito SRP auth. Hosted on S3 + CloudFront."],
    ]
    col_ws = [Inches(2.5), Inches(2.3), Inches(7.53)]
    tbl(sl, rows, CL, CY+Inches(0.1), CW, Inches(5.6), col_ws)

    # ── 7. PLACEHOLDER: Application Architecture ──────────────────────────────
    ph_slide(prs, "Application Architecture Diagram",
             "Show: Browser → CloudFront/ALB → Services → RDS + Redis + SQS + S3")

    # ── 8. EVENT-DRIVEN ARCHITECTURE ─────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Event-Driven Architecture")

    card(sl, CL, CY+Inches(0.1), C1W, Inches(5.7), "Synchronous Path")
    ml(sl, [
        (f"▶  Browser / SPA", WHITE),
        sub("Cognito SRP auth → JWT token"),
        (f"▶  AWS CloudFront → ALB", WHITE),
        sub("Wildcard ACM cert  *. rnld101.xyz"),
        (f"▶  appointment-service  |  report-service", WHITE),
        sub("FastAPI async, port 8000"),
        sub("PostgreSQL (RDS) via SQLAlchemy async"),
        (f"▶  Redis  (in-cluster, ephemeral)", WHITE),
        sub("5-min slot-locks for concurrent booking"),
        sub("JWKS public-key cache"),
    ], CL+Inches(0.15), CY+Inches(0.55), C1W-Inches(0.3), Inches(4.9),
       sz=10.5, spc=5)

    card(sl, C2L, CY+Inches(0.1), C2W, Inches(2.7), "Async / Event Path  (SQS)")
    ml(sl, [
        (f"▶  appointment-service publishes", WHITE),
        sub("appointment.booked  →  SQS queue"),
        sub("Decoupled: booking succeeds even if notification fails"),
        (f"▶  notification-service long-polls SQS", WHITE),
        sub("Receives event, sends email via SES"),
        sub("On error: logs + continues (DLQ handles retries)"),
        sub("Deletes message only on success"),
    ], C2L+Inches(0.15), CY+Inches(0.55), C2W-Inches(0.3), Inches(2.1), sz=10.5, spc=5)

    card(sl, C2L, CY+Inches(3.0), C2W, Inches(2.8), "AI Pipeline  (S3 → Lambda)")
    ml(sl, [
        (f"▶  PDF uploaded to S3 reports bucket", WHITE),
        sub("Triggers Lambda (S3 ObjectCreated event)"),
        (f"▶  Lambda: Textract OCR extracts full text", WHITE),
        (f"▶  Bedrock Nova Lite generates summary", WHITE),
        (f"▶  Titan Embeddings v2 vectorises chunks", WHITE),
        sub("Stored in pgvector (RDS Postgres)"),
        (f"▶  report-service RAG chat uses pgvector search", WHITE),
        sub("Document-scoped: no cross-patient data bleed"),
    ], C2L+Inches(0.15), CY+Inches(3.45), C2W-Inches(0.3), Inches(2.2), sz=10.5, spc=4)

    # ── 9. SECTION: CLOUD INFRA ───────────────────────────────────────────────
    sec_slide(prs, "02", "Cloud Infrastructure",
              "AWS Services  |  Multi-AZ VPC  |  EKS  |  RDS  |  CloudFront")

    # ── 10. AWS SERVICES USED ─────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "AWS Services Used")

    card(sl, CL, CY+Inches(0.05), C1W+Inches(0.1), Inches(2.85), "Compute & Networking")
    ml(sl, [
        bull("EKS v1.31  —  managed Kubernetes control plane"),
        bull("EC2  (c7i-flex.large)  —  managed node group + Karpenter nodes"),
        bull("VPC  —  public / private / DB subnets across 2 AZs,  NAT Gateway"),
        bull("ALB  —  provisioned by AWS Load Balancer Controller via Ingress"),
        bull("CloudFront  —  CDN for SPA  (OAC, HTTPS, wildcard ACM cert)"),
        bull("Lambda  —  AI processing pipeline  (SAM-managed)"),
        bull("Route53 + ACM  —  wildcard cert  *.rnld101.xyz"),
    ], CL+Inches(0.15), CY+Inches(0.5), C1W-Inches(0.15), Inches(2.2), sz=10.5, spc=4)

    card(sl, CL, CY+Inches(3.1), C1W+Inches(0.1), Inches(2.75), "Storage & Messaging")
    ml(sl, [
        bull("RDS PostgreSQL  (db.t4g.micro)  —  primary DB + pgvector"),
        bull("S3  —  reports bucket  (KMS-encrypted, versioned)  + frontend bucket"),
        bull("ECR  —  4 container image repos  (KMS-encrypted, immutable tags)"),
        bull("SQS  —  appointment.booked  event queue  (120s visibility timeout)"),
        bull("SES  —  outbound email  (domain identity, DKIM CNAMEs via Terraform)"),
    ], CL+Inches(0.15), CY+Inches(3.55), C1W-Inches(0.15), Inches(2.1), sz=10.5, spc=5)

    card(sl, C2L-Inches(0.1), CY+Inches(0.05), C2W+Inches(0.1), Inches(2.85), "Identity & Secrets")
    ml(sl, [
        bull("Cognito  —  user pool, SRP auth, patient / staff roles"),
        bull("Secrets Manager  —  DB URL, Grafana creds  (ESO-pulled at runtime)"),
        bull("SSM Parameter Store  —  non-sensitive config  (/lablumen/config/*)"),
        bull("KMS CMK  —  encrypts ECR images and Secrets Manager values"),
        bull("IAM + OIDC  —  GitHub OIDC federation  (4 roles, zero static keys)"),
        bull("IRSA  —  pod-level AWS identity via ServiceAccount annotation"),
    ], C2L+Inches(0.0), CY+Inches(0.5), C2W-Inches(0.15), Inches(2.2), sz=10.5, spc=4)

    card(sl, C2L-Inches(0.1), CY+Inches(3.1), C2W+Inches(0.1), Inches(2.75), "AI & Observability")
    ml(sl, [
        bull("Bedrock  (Nova Lite)  —  lab report AI summarisation"),
        bull("Textract  —  PDF / image text extraction for lab reports"),
        bull("EKS Control Plane Logs  →  CloudWatch"),
        bull("Prometheus  —  cluster metrics  (in-cluster, kube-prometheus-stack)"),
        bull("Grafana  —  dashboards at  grafana.rnld101.xyz  (ESO-managed creds)"),
        bull("External DNS  —  syncs K8s Ingress hosts  →  Route53 A records"),
    ], C2L+Inches(0.0), CY+Inches(3.55), C2W-Inches(0.15), Inches(2.1), sz=10.5, spc=4)

    # ── 11. PLACEHOLDER: AWS Architecture ────────────────────────────────────
    ph_slide(prs, "AWS Architecture Diagram",
             "Show: VPC → EKS (private) + RDS (DB subnet) | CloudFront → S3 | SQS | Lambda")

    # ── 12. SECTION: TERRAFORM ────────────────────────────────────────────────
    sec_slide(prs, "03", "Terraform — Infrastructure as Code",
              "12 Modules  |  Remote State  |  OIDC  |  IRSA  |  Infracost")

    # ── 13. TERRAFORM OVERVIEW ────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Terraform — How Infrastructure is Provisioned")

    ml(sl, [
        bull("Zero manual AWS console setup — every resource is Terraform-managed"),
        bull("Remote state: S3 bucket  lablumen-tfstate-<account-id>  with native S3 locking  (use_lockfile = true)"),
        bull("Bootstrap folder: creates state bucket with  prevent_destroy  before the main stack — separate concern"),
        bull("Account-portable: no hardcoded account IDs — derived via  data.aws_caller_identity"),
        bull("All resources tagged:  Environment  +  Owner  on every resource"),
        bull("terraform fmt  +  terraform validate  pass cleanly;  terraform.tfvars  drives all values"),
    ], CL, CY+Inches(0.05), C1W+Inches(0.5), Inches(3.2), sz=11, spc=7)

    card(sl, CL, CY+Inches(3.3), C1W+Inches(0.5), Inches(2.55), "4 GitHub OIDC IAM Roles — Zero Static Keys")
    ml(sl, [
        bull("lablumen-tf-plan     — read-only + state  (runs on every PR automatically)"),
        bull("lablumen-tf-apply    — admin  (runs on main only, after human approval)"),
        bull("lablumen-app-ci-ecr  — ECR push  (triggered on merge to each service repo)"),
        bull("lablumen-frontend-deploy  — S3 sync + CloudFront invalidation"),
    ], CL+Inches(0.15), CY+Inches(3.75), C1W+Inches(0.3), Inches(1.9), sz=10.5, spc=5)

    codebox(sl, [
        'terraform {',
        '  backend "s3" {',
        '    bucket       = "lablumen-tfstate-261523981519"',
        '    key          = "terraform.tfstate"',
        '    region       = "us-east-1"',
        '    use_lockfile = true   # native S3 locking (no DynamoDB)',
        '  }',
        '}',
    ], C2L+Inches(0.2), CY+Inches(0.05), C2W-Inches(0.2), Inches(2.05), "backend.tf")

    codebox(sl, [
        '# OIDC trust (no static credentials)',
        'resource "aws_iam_role" "app_ci_ecr" {',
        '  assume_role_policy = jsonencode({',
        '    Statement = [{',
        '      Principal = {',
        '        Federated = aws_iam_openid_connect_provider',
        '                     .github.arn',
        '      }',
        '      Condition = { StringLike = {',
        '        "token.actions.githubusercontent.com:sub" =',
        '          "repo:lablumen/*:ref:refs/heads/main"',
        '      }}',
        '    }]',
        '  })',
        '}',
    ], C2L+Inches(0.2), CY+Inches(2.25), C2W-Inches(0.2), Inches(3.65), "modules/iam — OIDC role")

    # ── 14. TERRAFORM MODULES ─────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Terraform — Module Structure")

    codebox(sl, [
        "modules/",
        "  vpc            VPC, 3 subnet tiers (public/private/DB), NAT GW",
        "  eks            EKS control plane, managed node group, Karpenter,",
        "                 Access Entries API (cluster-admin for tf-apply)",
        "  rds            Postgres db.t4g.micro, Secrets Manager creds",
        "  s3             Reports bucket (KMS, versioned) + SAM artifacts",
        "  cloudfront     OAC distribution + Route53 alias (toggleable)",
        "  ecr            4 image repos (KMS encrypted, immutable tags)",
        "  sqs            appointment.booked queue (120s visibility)",
        "  ses            Domain identity + DKIM CNAMEs in Route53",
        "  cognito        User pool + app client (SRP flow)",
        "  secretsmanager Secret shells (populated by ops, pulled by ESO)",
        "  ssm            Config params: /lablumen/config/*",
        "  iam            GitHub OIDC provider, 4 pipeline roles,",
        "                 all IRSA roles (ESO, report-svc, notify-svc,",
        "                 Karpenter, LBC, external-dns, ai-lambda)",
        "  lambda         AI function (SAM-managed, count-gated)",
    ], CL, CY+Inches(0.05), C1W+Inches(0.4), Inches(5.75), "modules/ layout")

    card(sl, C2L+Inches(0.1), CY+Inches(0.05), C2W-Inches(0.1), Inches(2.6), "Key Terraform Features")
    ml(sl, [
        bull("IRSA  —  pod-level AWS identity for 6 service accounts"),
        bull("EKS Access Entries API  —  cluster-admin grant for tf-apply"),
        bull("KMS CMK  —  shared platform key; encrypts ECR + Secrets Manager"),
        bull("Infracost  —  cost estimate on every plan, PR comment"),
        bull("Checkov  —  IaC security scan; SARIF  →  GitHub Security tab"),
        bull("enable_cloudfront toggle  —  count gating for new-account bring-up"),
    ], C2L+Inches(0.25), CY+Inches(0.5), C2W-Inches(0.35), Inches(2.0), sz=10.5, spc=5)

    card(sl, C2L+Inches(0.1), CY+Inches(2.85), C2W-Inches(0.1), Inches(2.85), "IRSA — Pod Identity (No Static Creds)")
    ml(sl, [
        bull("report-service  →  S3 (GetObject/PutObject) + Bedrock"),
        bull("notification-service  →  SQS (receive/delete) + SES (SendEmail)"),
        bull("ESO (lablumen-eso)  →  SecretsManager + SSM GetParameter"),
        bull("Karpenter  →  EC2 CreateFleet, DescribeInstances, etc."),
        bull("AWS LBC  →  ElasticLoadBalancing full access"),
        bull("external-dns  →  Route53 ChangeResourceRecordSets"),
    ], C2L+Inches(0.25), CY+Inches(3.3), C2W-Inches(0.35), Inches(2.2), sz=10.5, spc=5)

    # ── 15. SECTION: KUBERNETES ───────────────────────────────────────────────
    sec_slide(prs, "04", "Kubernetes on EKS",
              "Namespaces  |  Helm  |  ArgoCD  |  Karpenter  |  Grafana")

    # ── 16. CLUSTER LAYOUT ────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "EKS Cluster Layout")

    card(sl, CL, CY+Inches(0.05), C1W, Inches(3.1), "Namespaces")
    rows_ns = [
        ["Namespace", "Contents"],
        ["lablumen", "Production: appointment, report, notification, redis, frontend (2 replicas each)"],
        ["lablumen-dev", "Dev: same services (1 replica each), shorter resource limits"],
        ["argocd", "GitOps controller, app-of-apps bootstrap"],
        ["monitoring", "kube-prometheus-stack: Prometheus + Grafana + Alertmanager"],
        ["external-secrets", "ESO operator — pulls Secrets Manager + SSM"],
        ["kube-system", "Karpenter, AWS LBC, CoreDNS, external-dns, metrics-server"],
    ]
    tbl(sl, rows_ns, CL, CY+Inches(0.5), C1W, Inches(5.35),
        [Inches(1.9), Inches(3.85)])

    card(sl, C2L, CY+Inches(0.05), C2W, Inches(5.75), "Helm Charts & GitOps Values")
    ml(sl, [
        (f"▶  charts/microservice", WHITE),
        sub("Parametric chart for all 3 backend services"),
        sub("Deployment, Service, Ingress, HPA, PDB, ServiceAccount, ExternalSecret"),
        sub("TopologySpreadConstraints (zone + node)"),
        sub("Full securityContext hardening on every pod"),
        (f"▶  charts/redis", WHITE),
        sub("Hardened ephemeral Redis (slot-locks + JWKS cache)"),
        (f"▶  Value overlay pattern:", WHITE),
        sub("global-values.yaml  →  ECR registry URL"),
        sub("services/<svc>/values.yaml      →  common config"),
        sub("services/<svc>/values-dev.yaml  →  dev image SHA + sizing"),
        sub("services/<svc>/values-prod.yaml →  prod semver tag"),
        (f"▶  ArgoCD ApplicationSet (multi-source):", WHITE),
        sub("Source 1: global-values.yaml"),
        sub("Source 2: values.yaml + values-<env>.yaml overlay"),
        sub("Source 3: chart path (charts/microservice)"),
        (f"▶  16 ArgoCD apps  —  all Synced + Healthy", GREEN),
    ], C2L+Inches(0.15), CY+Inches(0.5), C2W-Inches(0.3), Inches(5.05), sz=10.5, spc=4)

    # ── 17. K8s SECURITY ─────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Kubernetes — Security & Production Hardening")

    codebox(sl, [
        "# Applied to EVERY pod in charts/microservice",
        "securityContext:",
        "  runAsNonRoot:              true",
        "  runAsUser:                 10001",
        "  fsGroup:                   10001",
        "  readOnlyRootFilesystem:    true",
        "  allowPrivilegeEscalation:  false",
        "  capabilities:",
        "    drop: [ALL]",
        "  seccompProfile:",
        "    type: RuntimeDefault",
        "automountServiceAccountToken: false",
        "",
        "# /tmp writable via emptyDir (readOnly rootFS)",
        "volumes:",
        "  - name: tmp",
        "    emptyDir: {}",
    ], CL, CY+Inches(0.05), C1W+Inches(0.1), Inches(4.1), "Pod SecurityContext")

    codebox(sl, [
        "# Health probes (verified in source code)",
        "livenessProbe:",
        "  httpGet: { path: /healthz, port: 8000 }",
        "  initialDelaySeconds: 10",
        "readinessProbe:",
        "  httpGet: { path: /readyz, port: 8000 }",
        "  initialDelaySeconds: 5",
        "",
        "# Resource limits",
        "resources:",
        "  requests: { cpu: 100m, memory: 128Mi }",
        "  limits:   { cpu: 500m, memory: 512Mi }",
    ], CL, CY+Inches(4.25), C1W+Inches(0.1), Inches(1.55), "Health Probes & Resource Limits")

    card(sl, C2L-Inches(0.1), CY+Inches(0.05), C2W+Inches(0.1), Inches(2.75), "Resilience")
    ml(sl, [
        bull("HPA  (CPU-based)  —  appointment-service + report-service"),
        bull("PodDisruptionBudget  maxUnavailable: 1  (safe rolling updates)"),
        bull("TopologySpreadConstraints  —  spread across AZs + nodes"),
        bull("2 replicas in production  (no single point of failure)"),
        bull("Startup probe on every pod  (avoids premature liveness kills)"),
    ], C2L+Inches(0.05), CY+Inches(0.5), C2W-Inches(0.2), Inches(2.1), sz=10.5, spc=5)

    card(sl, C2L-Inches(0.1), CY+Inches(3.0), C2W+Inches(0.1), Inches(2.85), "Secrets — Zero Hardcoding")
    ml(sl, [
        bull("External Secrets Operator pulls at runtime:"),
        sub("Secrets Manager  →  DATABASE_URL  (encrypted with KMS CMK)"),
        sub("SSM Parameter Store  →  all app config (SQS URL, S3 bucket, Cognito IDs,"),
        sub("   Bedrock model, CORS origins, SES sender, presigned-URL TTL)"),
        bull("Two ClusterSecretStores:"),
        sub("aws-secrets-manager   (SecretsManager source)"),
        sub("aws-parameter-store   (SSM source)"),
        bull("IRSA: ServiceAccount annotation  →  IAM role  (no node-level creds)"),
        bull("No secrets in Git — ever"),
    ], C2L+Inches(0.05), CY+Inches(3.45), C2W-Inches(0.2), Inches(2.25), sz=10.5, spc=4)

    # ── 18. ARGOCD GITOPS ─────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "ArgoCD — GitOps Deployment Model")

    card(sl, CL, CY+Inches(0.05), C1W+Inches(0.3), Inches(3.05), "App-of-Apps Pattern")
    ml(sl, [
        bull("bootstrap/root-app.yaml  →  bootstraps everything from a single Application"),
        sub("Includes: AppProject, platform addons, ApplicationSets"),
        bull("AppProject: allowed repos, destinations, namespace RBAC"),
        bull("ApplicationSets  (list generator + multi-source)"),
        sub("services-dev.yaml   →  deploys to lablumen-dev"),
        sub("services-prod.yaml  →  deploys to lablumen"),
        bull("Dev: auto-sync  |  Prod: promote by GitHub Release"),
    ], CL+Inches(0.15), CY+Inches(0.5), C1W+Inches(0.1), Inches(2.4), sz=10.5, spc=5)

    card(sl, CL, CY+Inches(3.3), C1W+Inches(0.3), Inches(2.55), "Sync Waves (ordering guarantee)")
    ml(sl, [
        (f"  Wave 0  —  Platform addons", TEAL),
        sub("ESO, metrics-server, Karpenter, LBC, external-dns, ArgoCD"),
        (f"  Wave 1  —  ClusterSecretStores", TEAL),
        sub("ESO must be ready before ExternalSecret resources sync"),
        (f"  Wave 2  —  Application services + Redis", TEAL),
        sub("appointment, report, notification, redis, frontend"),
    ], CL+Inches(0.15), CY+Inches(3.75), C1W+Inches(0.1), Inches(1.95), sz=10.5, spc=5)

    card(sl, C2L-Inches(0.3), CY+Inches(0.05), C2W+Inches(0.3), Inches(2.35), "Platform Addons managed by ArgoCD")
    ml(sl, [
        chk("AWS Load Balancer Controller  (creates ALBs from Ingress)"),
        chk("External DNS  (syncs Ingress hosts  →  Route53 A records)"),
        chk("External Secrets Operator  (ESO)"),
        chk("Karpenter + CRD Application  (v1, split chart)"),
        chk("Metrics Server  (enables HPA)"),
        chk("kube-prometheus-stack  (Grafana + Prometheus + Alertmanager)"),
        chk("ArgoCD  (self-managed via App-of-Apps)"),
    ], C2L-Inches(0.15), CY+Inches(0.5), C2W+Inches(0.1), Inches(1.75), sz=10.5, spc=3)

    card(sl, C2L-Inches(0.3), CY+Inches(2.55), C2W+Inches(0.3), Inches(3.3), "GitOps Deploy Flow")
    ml(sl, [
        (f"▶  feature/* branch  →  PR", WHITE),
        sub("Full security gate: lint + SAST + SCA + container scan"),
        (f"▶  Merge to main  →  CI builds  :sha  image", WHITE),
        sub("Trivy gate  →  ECR push  →  yq bumps values-dev.yaml"),
        (f"▶  ArgoCD detects git diff  →  auto-syncs DEV", TEAL),
        sub("lablumen-dev namespace updated in ~30 seconds"),
        (f"▶  Validate in DEV  →  GitHub Release  v1.2.0", WHITE),
        sub("CI retags  :sha  →  :v1.2.0  (manifest copy, no rebuild)"),
        sub("yq bumps values-prod.yaml  →  ArgoCD syncs PROD"),
        (f"▶  ArgoCD self-heals on any drift", GREEN),
    ], C2L-Inches(0.15), CY+Inches(3.0), C2W+Inches(0.1), Inches(2.7), sz=10.5, spc=4)

    # ── 19. PLACEHOLDER: ArgoCD Screenshot ───────────────────────────────────
    ph_slide(prs, "ArgoCD — Applications Dashboard",
             "Show: All 22+ apps Synced + Healthy, project lablumen, dev + prod")

    # ── 20. KARPENTER ────────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Karpenter — Node Auto-Provisioner")

    card(sl, CL, CY+Inches(0.05), C1W, Inches(3.3), "What Karpenter Does")
    ml(sl, [
        bull("Replaces Cluster Autoscaler  —  provisions nodes in ~30 seconds"),
        bull("Watches unschedulable pods and launches the right instance type"),
        bull("EC2NodeClass: discovers subnets + SGs by tag"),
        sub("karpenter.sh/discovery = lablumen-eks"),
        bull("NodePool: on-demand, c7i-flex.large, CPU safety cap 20 vCPUs"),
        bull("Consolidation: removes underutilised nodes automatically"),
        sub("WhenEmptyOrUnderutilized  |  consolidateAfter: 1m"),
        bull("IRSA-authenticated  —  no credentials on nodes"),
        bull("Terraform provisions IRSA via karpenter sub-module"),
    ], CL+Inches(0.15), CY+Inches(0.5), C1W-Inches(0.3), Inches(2.65), sz=10.5, spc=5)

    codebox(sl, [
        "apiVersion: karpenter.sh/v1",
        "kind: NodePool",
        "spec:",
        "  template:",
        "    spec:",
        "      nodeClassRef:",
        "        kind: EC2NodeClass",
        "        name: default",
        "      requirements:",
        "        - key: karpenter.sh/capacity-type",
        "          operator: In",
        "          values: [on-demand]",
        "        - key: node.kubernetes.io/instance-type",
        "          operator: In",
        "          values: [c7i-flex.large]",
        "  limits:",
        "    cpu: '20'  # cluster-wide safety cap",
        "  disruption:",
        "    consolidationPolicy: WhenEmptyOrUnderutilized",
        "    consolidateAfter: 1m",
    ], CL, CY+Inches(3.5), C1W, Inches(2.35), "platform/karpenter/nodepool.yaml")

    card(sl, C2L, CY+Inches(0.05), C2W, Inches(5.75), "Monitoring Stack  (kube-prometheus-stack)")
    ml(sl, [
        (f"▶  Prometheus  (internal)", WHITE),
        sub("Cluster-wide metrics scraping"),
        sub("2-day retention, ephemeral storage"),
        sub("Internal only: kubectl port-forward for direct access"),
        (f"▶  Grafana  (external)", WHITE),
        sub("Exposed at  grafana.rnld101.xyz  via ALB + HTTPS"),
        sub("Admin credentials pulled from Secrets Manager via ESO"),
        sub("K8S Dashboard  —  33 workloads, 44 pods, 2 nodes visible"),
        (f"▶  Alertmanager  (included, ephemeral)", WHITE),
        (f"▶  ArgoCD UI  (external)", WHITE),
        sub("Exposed at  argocd.rnld101.xyz  via same ALB ingress group"),
        (f"▶  Cost-optimised: both share one ALB  (Ingress Group)", GREEN),
        sub("External DNS auto-creates Route53 A records for both domains"),
        (f"▶  kube-state-metrics + node-exporter  included", WHITE),
        sub("Full node CPU/memory/disk visibility in Grafana"),
    ], C2L+Inches(0.15), CY+Inches(0.5), C2W-Inches(0.3), Inches(5.1), sz=10.5, spc=4)

    # ── 21. PLACEHOLDER: Grafana ──────────────────────────────────────────────
    ph_slide(prs, "Grafana — K8s Dashboard",
             "Show: Node Resource Overview — CPU/Memory, Pod count, Network")

    # ── 22. SECTION: GITHUB ACTIONS ───────────────────────────────────────────
    sec_slide(prs, "05", "GitHub Actions — CI/CD",
              "Trunk-Based  |  Security Gates  |  GitOps Write-back  |  Prod Promotion")

    # ── 23. PIPELINE ARCHITECTURE + BRANCHING ────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "GitHub Actions — Architecture & Branching Strategy")

    card(sl, CL, CY+Inches(0.05), C1W, Inches(2.9), "Polyrepo Structure")
    ml(sl, [
        bull("5 service repos  —  one per microservice / frontend"),
        bull("1 shared repo  (lablumen-shared)  —  reusable workflow engine"),
        sub("Called by all service repos via  uses: lablumen/lablumen-shared/...@main"),
        sub("Single source of CI truth  —  no logic duplication"),
        bull("lablumen-k8s  —  GitOps values repo (image tags, config)"),
        bull("lablumen-terraform  —  IaC repo (separate pipeline)"),
    ], CL+Inches(0.15), CY+Inches(0.5), C1W-Inches(0.3), Inches(2.25), sz=10.5, spc=5)

    card(sl, CL, CY+Inches(3.15), C1W, Inches(2.7), "Trunk-Based Development")
    ml(sl, [
        bull("All work on short-lived  feature/*  branches"),
        bull("Direct push to  main  is BLOCKED  by branch protection rules"),
        bull("PRs required to merge  —  triggers full 4-gate security pipeline"),
        bull("Merge to main  →  automatic build + dev deploy"),
        bull("Short-lived branches, fast integration, no merge conflicts"),
        bull("GitHub Release  (semver)  →  production promotion"),
    ], CL+Inches(0.15), CY+Inches(3.6), C1W-Inches(0.3), Inches(2.1), sz=10.5, spc=5)

    card(sl, C2L, CY+Inches(0.05), C2W, Inches(5.75), "3 Reusable Workflow Templates  (lablumen-shared)")
    ml(sl, [
        (f"▶  service-pr.yml  —  PR security gate", TEAL),
        sub("Ruff lint + Pytest  (hard-fail)"),
        sub("SonarCloud SAST + quality gate  (blocks on fail)"),
        sub("Snyk SCA  (severity-threshold = high, hard-fail)"),
        sub("Trivy container scan  (CRITICAL/HIGH, image never pushed)"),
        ("", LGREY),
        (f"▶  service-build-push.yml  —  merge → dev deploy", TEAL),
        sub("Docker multistage build  →  Trivy gate  →  ECR push  (OIDC)"),
        sub("yq bumps  values-dev.yaml  in lablumen-k8s"),
        sub("git pull --rebase  retry loop  (race-safe for concurrent services)"),
        sub("ArgoCD detects diff  →  auto-syncs lablumen-dev  (~30s)"),
        ("", LGREY),
        (f"▶  service-release.yml  —  GitHub Release → prod", TEAL),
        sub("ECR manifest copy:  :sha  →  :v1.2.0  (no rebuild, immutable safe)"),
        sub("yq bumps  values-prod.yaml  →  ArgoCD syncs lablumen (prod)"),
        sub("Build-once / promote-by-retag:  exact binary from dev reaches prod"),
    ], C2L+Inches(0.15), CY+Inches(0.5), C2W-Inches(0.3), Inches(5.05), sz=10.5, spc=3)

    # ── 24. PR GATE ───────────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "PR Security Gate  (service-pr.yml)", "All 4 jobs are hard-fail — code cannot reach main without passing all simultaneously")

    # Flow boxes
    BW = Inches(2.6); BH2 = Inches(1.65); BY = CY + Inches(0.5)
    BPOSITIONS = [CL, CL+Inches(3.35), CL+Inches(6.7), CL+Inches(10.05)]
    BLABELS = ["lint-and-test", "sast", "sca", "container-scan"]
    BICONS  = ["✎", "⛔", "⛔", "⛔"]  # pencil, block symbols
    BTITLES = ["Ruff + Pytest", "SonarCloud\nSAST", "Snyk\nSCA", "Trivy\nContainer Scan"]
    BSUBS   = [
        "Python linter\n+ unit tests\n(hard-fail)",
        "Full static\nanalysis +\nquality gate wait",
        "Dependency\nvuln scan\nhigh threshold",
        "Builds image\nlocally, never\npushed to ECR",
    ]
    for i in range(4):
        rect(sl, BPOSITIONS[i], BY, BW, BH2, CH, TH, bw=1.0)
        rect(sl, BPOSITIONS[i], BY, BW, Inches(0.3), TH)
        tx(sl, BLABELS[i], BPOSITIONS[i]+Inches(0.1), BY+Inches(0.05),
           BW-Inches(0.2), Inches(0.25), sz=9, bold=True, col=WHITE)
        tx(sl, BTITLES[i], BPOSITIONS[i]+Inches(0.1), BY+Inches(0.35),
           BW-Inches(0.2), Inches(0.55), sz=14, bold=True, col=TEAL)
        tx(sl, BSUBS[i], BPOSITIONS[i]+Inches(0.1), BY+Inches(0.9),
           BW-Inches(0.2), Inches(0.75), sz=9.5, col=LGREY)
        if i < 3:
            tx(sl, "→", BPOSITIONS[i]+BW+Inches(0.05), BY+Inches(0.65),
               Inches(0.3), Inches(0.4), sz=20, bold=True, col=TEAL)

    tx(sl, "Triggered on every pull_request event across all 5 service repositories",
       CL, BY+BH2+Inches(0.2), CW, Inches(0.35), sz=11, col=LGREY, italic=True)

    card(sl, CL, BY+BH2+Inches(0.7), CW/2-Inches(0.1), Inches(2.7), "What each gate catches")
    ml(sl, [
        (f"  Ruff", TEAL), sub("Python code style violations, unused imports, type errors"),
        (f"  Pytest", TEAL), sub("Unit test failures, regressions in business logic"),
        (f"  SonarCloud SAST", TEAL), sub("SQL injection, hardcoded secrets, code smells, coverage"),
        (f"  Snyk SCA", TEAL), sub("Known CVEs in pip dependencies  (high severity threshold)"),
        (f"  Trivy", TEAL), sub("OS + library CVEs in the container image  (CRITICAL/HIGH)"),
    ], CL+Inches(0.15), BY+BH2+Inches(1.15), CW/2-Inches(0.3), Inches(2.1), sz=10, spc=3)

    card(sl, CL+CW/2+Inches(0.1), BY+BH2+Inches(0.7), CW/2-Inches(0.1), Inches(2.7), "Why this order matters")
    ml(sl, [
        bull("lint-and-test  runs first  —  fast feedback, cheap"),
        bull("sast, sca, container-scan  run in PARALLEL after lint"),
        sub("All 3 must pass for merge to unblock"),
        bull("Container image is built locally and scanned but NEVER pushed to ECR"),
        sub("Push only happens after merge to main (build-push workflow)"),
        bull("Quality gate  (sonar)  waits for analysis completion before returning result"),
        bull("No credentials needed for PR gate  —  zero AWS access required"),
    ], CL+CW/2+Inches(0.25), BY+BH2+Inches(1.15), CW/2-Inches(0.35), Inches(2.1), sz=10, spc=4)

    # ── 25. PLACEHOLDER: Pipeline Screenshot ─────────────────────────────────
    ph_slide(prs, "CI Pipeline Run Screenshot",
             "Show: GitHub Actions run with pr / deploy-dev jobs, status Success")

    # ── 26. PLACEHOLDER: SonarCloud ──────────────────────────────────────────
    ph_slide(prs, "SonarCloud SAST Screenshot",
             "Show: Project analysis, quality gate PASSED, code smells / coverage")

    # ── 27. DEV DEPLOY + PROD RELEASE ────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "CI/CD — Dev Deploy & Production Release")

    card(sl, CL, CY+Inches(0.05), C1W+Inches(0.2), Inches(5.75), "Dev Deploy  (service-build-push.yml)  —  merge to main")
    ml(sl, [
        (f"1.  Checkout  +  compute 7-char git SHA tag", WHITE),
        sub("Short SHA used as immutable image tag in ECR"),
        (f"2.  Configure AWS via GitHub OIDC", WHITE),
        sub("Assumes  lablumen-app-ci-ecr  role  —  no static keys"),
        (f"3.  Docker multistage build", WHITE),
        sub("builder stage  →  /opt/venv  |  runtime stage  →  non-root uid 10001"),
        (f"4.  Trivy container scan  (CRITICAL/HIGH = abort)", WHITE),
        sub("Runs against built image before any push to ECR"),
        (f"5.  Push  lablumen/<svc>:<sha>  to ECR", WHITE),
        sub("KMS-encrypted ECR, immutable tags, SHA is unique per commit"),
        (f"6.  GitOps write-back  to lablumen-k8s", WHITE),
        sub("yq  sets  .image.tag = <sha>  in  services/<svc>/values-dev.yaml"),
        sub("git pull --rebase  retry loop  (up to 5 retries)  —  race-safe"),
        sub("for concurrent pushes from multiple service repos"),
        (f"7.  ArgoCD detects diff  →  auto-syncs  lablumen-dev", GREEN),
        sub("New pod with updated image running in ~30 seconds"),
    ], CL+Inches(0.15), CY+Inches(0.5), C1W, Inches(5.05), sz=10, spc=4)

    card(sl, C2L-Inches(0.1), CY+Inches(0.05), C2W+Inches(0.1), Inches(5.75), "Production Release  (service-release.yml)")
    ml(sl, [
        (f"▶  Trigger: GitHub Release published  (e.g.  v1.2.0)", ORANGE),
        sub("Release must be created on the commit already in DEV"),
        (f"▶  Build-once / Promote-by-retag  strategy", TEAL),
        sub("Does NOT rebuild the image"),
        sub("The exact binary that passed all tests in DEV reaches PROD"),
        (f"▶  ECR manifest copy  (retag):", WHITE),
        sub("batch-get-image  →  get manifest for  :sha"),
        sub("put-image  →  write  :v1.2.0  (new tag, idempotent)"),
        sub("Works with IMMUTABLE ECR repos (semver tag is new)"),
        sub("No layer re-upload  —  just a new tag pointer"),
        (f"▶  GitOps write-back  to lablumen-k8s:", WHITE),
        sub("yq  sets  .image.tag = v1.2.0  in  values-prod.yaml"),
        sub("rebase-retry loop for concurrent service releases"),
        (f"▶  ArgoCD syncs  lablumen  (prod)  namespace", GREEN),
        sub("Zero-downtime rolling update, PDB protects availability"),
    ], C2L+Inches(0.05), CY+Inches(0.5), C2W-Inches(0.15), Inches(5.05), sz=10, spc=4)

    # ── 28. TERRAFORM PIPELINE ────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Terraform Pipeline", "Triggered on *.tf / *.tfvars changes  or  workflow_dispatch")

    # 4 pipeline stage boxes
    STAGE_W = Inches(2.85); STAGE_H = Inches(2.6)
    STAGE_X = [CL, CL+Inches(3.3), CL+Inches(6.6), CL+Inches(9.9)]
    STAGE_FILLS = [CH, CH, H(50, 30, 10), H(18, 50, 28)]
    STAGE_BORDERS = [TH, TH, OH, GH]
    STAGE_NUMS = ["1", "2", "3", "4"]
    STAGE_TITLES = ["Checkov\nIaC Scan", "Terraform Plan\n+ Infracost", "Manual\nApproval Gate", "Terraform\nApply"]
    STAGE_DETAILS = [
        "Scans all .tf files\nfor security issues\n\nUploads SARIF to\nGitHub Security tab\n\nSoft-fail (reports,\ndoes not block now)",
        "tf-plan OIDC role\n(read-only + state)\n\nfmt-check + validate\n+ plan -out tfplan\n\nInfracost estimate\nposted as PR comment",
        "GitHub Environment\n'production'\n\nRequired reviewers\nmust approve\n\nRuns on main only,\nnot on PRs",
        "tf-apply OIDC role\n(admin permissions)\n\nDownloads saved plan\napplies exact plan\n(no drift between\nreview and apply)",
    ]
    for i in range(4):
        rect(sl, STAGE_X[i], CY+Inches(0.15), STAGE_W, STAGE_H,
             STAGE_FILLS[i], STAGE_BORDERS[i], bw=1.5)
        tx(sl, STAGE_NUMS[i], STAGE_X[i]+Inches(0.12), CY+Inches(0.22),
           Inches(0.5), Inches(0.55), sz=24, bold=True,
           col=TEAL if i < 2 else (ORANGE if i == 2 else GREEN))
        tx(sl, STAGE_TITLES[i], STAGE_X[i]+Inches(0.12), CY+Inches(0.75),
           STAGE_W-Inches(0.25), Inches(0.7), sz=13, bold=True, col=WHITE)
        tx(sl, STAGE_DETAILS[i], STAGE_X[i]+Inches(0.12), CY+Inches(1.45),
           STAGE_W-Inches(0.25), Inches(1.2), sz=9, col=LGREY)
        if i < 3:
            tx(sl, "→", STAGE_X[i]+STAGE_W+Inches(0.08), CY+Inches(1.1),
               Inches(0.28), Inches(0.5), sz=22, bold=True,
               col=(ORANGE if i == 1 else TEAL))

    card(sl, CL, CY+Inches(2.9), CW/2-Inches(0.1), Inches(2.9), "Two separate OIDC roles by design")
    ml(sl, [
        bull("tf-plan  role  =  read-only + state bucket access"),
        sub("Can read all resources, but cannot create/modify/delete"),
        sub("Runs automatically on every PR and push to main"),
        sub("Safe to run without human oversight"),
        bull("tf-apply  role  =  full admin + EKS cluster-admin"),
        sub("Only runs after human approval via 'production' Environment"),
        sub("Applies the SAVED plan artifact (no re-plan on apply)"),
        sub("EKS cluster-admin needed for kubernetes_* resources"),
    ], CL+Inches(0.15), CY+Inches(3.35), CW/2-Inches(0.3), Inches(2.3), sz=10, spc=5)

    card(sl, CL+CW/2+Inches(0.1), CY+Inches(2.9), CW/2-Inches(0.1), Inches(2.9), "One-Click Destroy Workflow")
    ml(sl, [
        bull("terraform-destroy.yml  (workflow_dispatch  with confirm='destroy')"),
        bull("2-phase approach to avoid orphaned AWS resources:"),
        sub("Phase 1: scale down ArgoCD app-controller  →  kubectl delete"),
        sub("         all Ingress + LB services  →  wait for ALBs to drain"),
        sub("         (LBC-managed ALBs would block VPC destroy otherwise)"),
        sub("Phase 2: terraform destroy  (VPC tears down cleanly)"),
        bull("Gated by GitHub  'production'  Environment  (required reviewers)"),
        bull("ECR module: force_delete=true  (empties repos on destroy)"),
    ], CL+CW/2+Inches(0.25), CY+Inches(3.35), CW/2-Inches(0.35), Inches(2.3), sz=10, spc=4)

    # ── 29. PLACEHOLDER: Infracost ────────────────────────────────────────────
    ph_slide(prs, "Infracost Cost Estimate",
             "Show: PR comment with monthly cost estimate per resource")

    # ── 30. SECTION: SECURITY ────────────────────────────────────────────────
    sec_slide(prs, "06", "Security & Results",
              "Zero static credentials end-to-end  |  Platform achievements")

    # ── 31. SECURITY SUMMARY ─────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Security — End to End")

    card(sl, CL, CY+Inches(0.05), C1W+Inches(0.2), Inches(2.75), "Zero Static Credentials")
    ml(sl, [
        chk("CI authenticates to AWS via GitHub OIDC  (no access keys in secrets)"),
        chk("Pods authenticate to AWS via IRSA  (ServiceAccount-bound IAM roles)"),
        chk("Config + secrets pulled at runtime by ESO  (Secrets Manager + SSM)"),
        chk("No secrets in Git — Trivy and SonarCloud enforce this on every PR"),
        chk("KMS CMK encrypts ECR image layers and Secrets Manager values"),
    ], CL+Inches(0.15), CY+Inches(0.5), C1W, Inches(2.05), sz=10.5, spc=5)

    card(sl, CL, CY+Inches(2.95), C1W+Inches(0.2), Inches(2.85), "Code & Container Security")
    ml(sl, [
        chk("SAST: SonarCloud quality gate blocks PR merge on findings"),
        chk("SCA: Snyk hard-fails on high-severity dependency vulnerabilities"),
        chk("Container: Trivy  CRITICAL/HIGH  blocks PR gate and pre-ECR-push"),
        chk("IaC: Checkov SARIF uploaded to GitHub Security tab on every push"),
        chk("Multistage Dockerfiles  —  slim runtime image, non-root UID 10001"),
        chk("Dependency upgrades in runtime stage  (apt upgrade)  clear fixable CVEs"),
    ], CL+Inches(0.15), CY+Inches(3.4), C1W, Inches(2.2), sz=10.5, spc=4)

    card(sl, C2L-Inches(0.1), CY+Inches(0.05), C2W+Inches(0.1), Inches(2.75), "Runtime & Network Security")
    ml(sl, [
        chk("readOnlyRootFilesystem + drop ALL capabilities on every pod"),
        chk("seccompProfile: RuntimeDefault  (kernel syscall filtering)"),
        chk("RDS in isolated database subnets  —  no public endpoint"),
        chk("S3 reports bucket: KMS-encrypted, versioned, private  (no public access)"),
        chk("Direct push to main blocked  —  branch protection + required PRs"),
        chk("automountServiceAccountToken: false  on all pods"),
    ], C2L+Inches(0.05), CY+Inches(0.5), C2W-Inches(0.2), Inches(2.05), sz=10.5, spc=4)

    card(sl, C2L-Inches(0.1), CY+Inches(2.95), C2W+Inches(0.1), Inches(2.85), "Security Gate Summary")
    ml(sl, [
        chk("No hardcoded credentials in any file or git history"),
        chk("Container does not run as root  (UID 10001)"),
        chk("Multistage Docker builds  (separate builder + runtime)"),
        chk("Container image scanned before PR merge and before ECR push"),
        chk("Secrets in Kubernetes via ExternalSecret  (not ConfigMaps or env literals)"),
        chk("RBAC: ServiceAccount has minimal IAM permissions  (least-privilege IRSA)"),
        chk(".gitignore excludes *.tfstate, .env, credentials files"),
    ], C2L+Inches(0.05), CY+Inches(3.4), C2W-Inches(0.2), Inches(2.2), sz=10.5, spc=3)

    # ── 32. RESULTS ───────────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    hbar(sl, "Platform Achievements")

    card(sl, CL, CY+Inches(0.05), C1W+Inches(0.1), Inches(2.75), "Infrastructure")
    ml(sl, [
        chk("Full AWS stack  —  12 Terraform modules, remote S3 state"),
        chk("EKS v1.31 with Karpenter auto-scaling across 2 AZs"),
        chk("Multi-AZ VPC, RDS in isolated subnets, CloudFront SPA"),
        chk("OIDC + IRSA  —  zero static credentials end-to-end"),
        chk("All resources tagged: Environment + Owner"),
    ], CL+Inches(0.15), CY+Inches(0.5), C1W-Inches(0.1), Inches(2.1), sz=10.5, spc=4)

    card(sl, CL, CY+Inches(3.0), C1W+Inches(0.1), Inches(2.85), "Application")
    ml(sl, [
        chk("4 backend microservices + React frontend  (all containerised)"),
        chk("Event-driven: SQS decouples appointment from notification"),
        chk("AI pipeline: Textract + Bedrock + pgvector  (Lambda/SAM)"),
        chk("Cognito auth: separate patient + staff RBAC roles"),
        chk("RAG chat on lab reports  —  document-scoped pgvector search"),
    ], CL+Inches(0.15), CY+Inches(3.45), C1W-Inches(0.1), Inches(2.2), sz=10.5, spc=4)

    card(sl, C2L-Inches(0.1), CY+Inches(0.05), C2W+Inches(0.1), Inches(2.75), "DevOps")
    ml(sl, [
        chk("Trunk-based + branch protection  (PR-only merges to main)"),
        chk("4-gate security pipeline on every PR  (SAST+SCA+Trivy+Ruff)"),
        chk("Build-once / promote-by-retag CD  (exact binary dev → prod)"),
        chk("GitOps: ArgoCD App-of-Apps, 16 apps Synced + Healthy"),
        chk("Prometheus + Grafana observability  (kube-prometheus-stack)"),
        chk("One-click destroy pipeline with 2-phase LBC cleanup"),
    ], C2L+Inches(0.05), CY+Inches(0.5), C2W-Inches(0.2), Inches(2.1), sz=10.5, spc=4)

    card(sl, C2L-Inches(0.1), CY+Inches(3.0), C2W+Inches(0.1), Inches(2.85),
         "Bonus Pillars Achieved  (each = +5%)")
    ml(sl, [
        (f"★  Event-Driven Architecture  —  SQS (appointment → notification)", ORANGE),
        (f"★  AI / ML Integration  —  Bedrock + Textract + pgvector", ORANGE),
        (f"★  Advanced Monitoring  —  Prometheus + Grafana + Alertmanager", ORANGE),
        (f"★  GitOps Deployment  —  ArgoCD  (App-of-Apps, 16 apps)", ORANGE),
        (f"★  Cost Estimation  —  Infracost on every terraform plan", ORANGE),
    ], C2L+Inches(0.05), CY+Inches(3.45), C2W-Inches(0.2), Inches(2.2), sz=11, spc=5)

    # ── 33. THANK YOU ────────────────────────────────────────────────────────
    sl = ns(prs); set_bg(sl, 10, 18, 52)
    rect(sl, 0, Inches(1.8), SW, Inches(4.0), H(14, 26, 70))
    rect(sl, 0, Inches(1.8), Inches(0.1), Inches(4.0), TH)
    rect(sl, 0, SH-Inches(0.07), SW, Inches(0.07), OH)
    tx(sl, "Thank You", Inches(1), Inches(2.0), Inches(8), Inches(1.4),
       sz=60, bold=True, col=WHITE)
    tx(sl, "Questions?", Inches(1), Inches(3.4), Inches(8), Inches(0.8),
       sz=28, col=TEAL)
    rect(sl, Inches(1), Inches(4.35), Inches(4), Inches(0.035), TH)

    card(sl, Inches(7.5), Inches(2.1), Inches(5.4), Inches(2.05), "Repositories")
    ml(sl, [
        ("lablumen/lablumen-terraform", LTEAL),
        ("lablumen/lablumen-k8s", LTEAL),
        ("lablumen/lablumen-shared", LTEAL),
        ("lablumen/lablumen-*-service", LTEAL),
    ], Inches(7.65), Inches(2.5), Inches(5.1), Inches(1.5),
       sz=10, fn="Courier New", spc=5)

    card(sl, Inches(7.5), Inches(4.3), Inches(5.4), Inches(2.35), "Live Endpoints")
    ml(sl, [
        ("api.rnld101.xyz          (prod API)", LTEAL),
        ("api-dev.rnld101.xyz   (dev API)", LTEAL),
        ("argocd.rnld101.xyz    (ArgoCD UI)", LTEAL),
        ("grafana.rnld101.xyz  (Grafana)", LTEAL),
        ("app.rnld101.xyz         (Frontend SPA)", LTEAL),
    ], Inches(7.65), Inches(4.7), Inches(5.1), Inches(1.8),
       sz=10, fn="Courier New", spc=5)

    tx(sl, "LabLumen  |  AWS EKS  |  GitOps  |  GitHub Actions  |  June 2026",
       Inches(1), Inches(4.55), Inches(6), Inches(0.45),
       sz=10, col=RGBColor(100, 125, 175), italic=True)

    # ── SAVE ──────────────────────────────────────────────────────────────────
    out = os.path.join(os.path.dirname(__file__), "lablumen-presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    build()
